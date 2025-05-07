from typing import Any
import google.generativeai as genai
import os
import mimetypes
from .logger import Logger
from pydantic import BaseModel


class Slide(BaseModel):
    title: str
    content: list[str]


class AiServices:
    def __init__(self, api_key=None, model_name="gemini-2.0-flash", logger=None,):
        """Initialize AiServices with API key and default model"""
        self.api_key = api_key 
        self.model_name = "gemini-2.0-flash"
      
        # Configure Gemini API
        genai.configure(api_key=self.api_key)
        
        # Initialize logger
        self.logger = logger
        
        # Default generation config
        self.generation_config = {
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 40,
            "max_output_tokens": 8192,
            "response_mime_type": "text/plain",
        }
        
        self.logger.info(f"AiServices initialized with model: {self.model_name}")

    def upload_to_gemini(self, path: str, mime_type: str = None):
        """Uploads the given file to Gemini with automatic mime type detection if not provided."""
        self.logger.info(f"Started uploading: {path}")
        if not os.path.exists(path):
            error_msg = f"File not found at path: {path}"
            self.logger.error(error_msg)
            raise FileNotFoundError(error_msg)

        # Normalize the path to handle spaces and special characters
        path = os.path.normpath(path)
        
        # Auto-detect mime type if not provided
        if mime_type is None:
            mime_type = mimetypes.guess_type(path)[0]
            if mime_type is None:
                self.logger.warning(f"Could not detect mime type for {path}. Gemini will attempt to determine the type.")

        try:
            file = genai.upload_file(str(path), mime_type=mime_type)
            self.logger.info(f"Uploaded file '{file.display_name}' as: {file.uri} (type: {mime_type or 'auto-detected'})")
            return file
        except Exception as e:
            error_msg = f"Upload failed: {str(e)}"
            self.logger.error(error_msg)
            raise

    def save_response(self, content, output_path, output_type=None):
        """
        Save response content in the appropriate format based on output_type or file extension
        
        Args:
            content (str): Content to save
            output_path (str): Path to save the file
            output_type (str, optional): Force specific output type (pdf, txt, docx)
        """
        # Determine the output type from path extension if not explicitly provided
        if output_type is None:
            _, ext = os.path.splitext(output_path)
            output_type = ext[1:].lower() if ext else "txt"
        
        self.logger.info(f"Saving response as {output_type} file to {output_path}")
        
        try:
            
            if output_type == "txt":
                content = content.replace("*", " ")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
            
            elif output_type == "pdf":
                try:
                    content = content.replace("*", " ")

                    from reportlab.lib.pagesizes import letter
                    from reportlab.platypus import SimpleDocTemplate, Paragraph
                    from reportlab.lib.styles import getSampleStyleSheet
                    
                    doc = SimpleDocTemplate(output_path, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = [Paragraph(line, styles["Normal"]) for line in content.split('\n')]
                    doc.build(story)
                except ImportError:
                    self.logger.error("reportlab library not found. Install with: pip install reportlab")
                    raise
                    
            elif output_type == "docx":
                try:
                    content = content.replace("*", " ")
                    
                    from docx import Document
                    
                    document = Document()
                    for paragraph in content.split('\n'):
                        if paragraph.strip():  # Skip empty paragraphs
                            document.add_paragraph(paragraph)
                    document.save(output_path)
                except ImportError:
                    self.logger.error("python-docx library not found. Install with: pip install python-docx")
                    raise
            elif output_type == "pptx":
                self.create_presentation_from_json(content, output_path)
            else:
                # Default to text file for unknown types
                self.logger.warning(f"Unknown output type '{output_type}', saving as text file")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
            
            self.logger.info(f"Response saved to {output_path}")
            return output_path
            
        except Exception as e:
            error_msg = f"Error saving {output_type} file: {str(e)}"
            self.logger.error(error_msg)
            raise

    def process_mixed_files(self, file_paths: list, system_prompt: str, user_prompt: str, 
                           output_path: str = "response.txt", model_name: str = None, output_type: str = None):
        """Process multiple files of mixed types with Gemini using chat mode"""
        
        # Determine the output type from path extension if not explicitly provided
        if output_type is None:
            _, ext = os.path.splitext(output_path)
            output_type = ext[1:].lower() if ext else "txt"
        
        # Use provided model name or default to instance model name
        model_name = model_name or self.model_name
        
        # Set up specific configuration for pptx output
        if output_type == "pptx":
            # Use JSON response format for presentations
            temp_config = self.generation_config.copy()
            temp_config["response_mime_type"] = "application/json"
            
            # Use client for schema-based generation
            client = genai.GenerativeModel(
                model_name=model_name,
                generation_config=temp_config,
                system_instruction="List slides from the attached content in JSON FORMAT with title and content for each slide."
            )
        else:
            # Initialize model with system prompt for non-pptx outputs
            client = genai.GenerativeModel(
                model_name=model_name,
                generation_config=self.generation_config,
                system_instruction=system_prompt
            )
        
        try:
            # Upload all files to Gemini
            uploaded_files = []
            for file_path in file_paths:
                self.logger.info(f"UPLOADING: {file_path}")
                uploaded_file = self.upload_to_gemini(file_path)
                uploaded_files.append(uploaded_file)
            
            # Start chat session
            chat = client.start_chat()
            
            # Create message content with all files and the prompt
            if output_type == "pptx":
                user_message = "Create a presentation from the content in the uploaded files. Format the response as JSON with slides object containing key-value pairs for each slide."
                message_content = uploaded_files + [user_message]
            else:
                message_content = uploaded_files + [user_prompt]
            
            # Send message with files and prompt
            self.logger.info("Sending message to Gemini with files and prompt")
            response = chat.send_message(
                message_content,
                request_options={"timeout": 1000},
            )
            
            self.logger.info("Received response from Gemini")
            self.logger.info(response.text)
            
            # Save response in the appropriate format
            self.save_response(response.text, output_path, output_type)
            
            return response.text
            
        except Exception as e:
            error_msg = f"Error processing files: {str(e)}"
            self.logger.error(error_msg)
            raise
    
    def set_generation_config(self, temperature=None, top_p=None, top_k=None, 
                             max_output_tokens=None, response_mime_type=None):
        """Update generation configuration parameters"""
        if temperature is not None:
            self.generation_config["temperature"] = temperature
        if top_p is not None:
            self.generation_config["top_p"] = top_p
        if top_k is not None:
            self.generation_config["top_k"] = top_k
        if max_output_tokens is not None:
            self.generation_config["max_output_tokens"] = max_output_tokens
        if response_mime_type is not None:
            self.generation_config["response_mime_type"] = response_mime_type
        
        self.logger.info(f"Updated generation config: {self.generation_config}")
        return self.generation_config

    def create_presentation_from_json(self, slide_data, output_path):
        """
        Create a PowerPoint presentation from JSON slide data.
        
        Args:
            slide_data (dict or str): Dictionary or JSON string containing slide data
            output_path (str): Path to save the presentation
        """
        import json
        from pptx import Presentation
        from pptx.util import Inches
        
        # Handle string input (parse JSON)
        if isinstance(slide_data, str):
            try:
                slide_data = json.loads(slide_data)
            except json.JSONDecodeError as e:
                self.logger.error(f"Failed to parse JSON: {e}")
                raise
        
        # Create a presentation object
        prs = Presentation()
        self.logger.info(f"Slide data: {slide_data}")
        # Check if slides is a list (new format) or dict (old format)
        if isinstance(slide_data, list):
            # Process direct list of slides format
            for slide_info in slide_data:
                title = slide_info.get("title", "")
                content = slide_info.get("content", "")
                
                # Use title and content slide layout
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Set the title
                title_shape = slide.shapes.title
                title_shape.text = title
                
                # Set the content
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = content
        elif isinstance(slide_data.get("slides"), list):
            # Process list-based slides format
            for slide_info in slide_data["slides"]:
                title = slide_info.get("title", "")
                content = slide_info.get("content", "")
                
                # Use title and content slide layout
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Set the title
                title_shape = slide.shapes.title
                title_shape.text = title
                
                # Set the content
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = content

        elif isinstance(slide_data.get("slides"), dict):
            # Process original dict-based slides format
            for slide_id, slide_info in slide_data["slides"].items():
                title = slide_info.get("title", "")
                content = slide_info.get("content", "")
                
                # Use title and content slide layout
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Set the title
                title_shape = slide.shapes.title
                title_shape.text = title
                
                # Set the content
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = content
        else:
            self.logger.error(f"Invalid slide data format: {slide_data}")
            raise ValueError("Slide data missing or in incorrect format")
            
        # Save the presentation
        prs.save(output_path)
        return output_path